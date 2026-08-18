#!/usr/bin/env python3
"""Build BarathX cold-outreach PPTX (no fee slides)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "BarathX-Creator-Collab-Brief.pptx"

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x4A, 0x55, 0x63)
SOFT = RGBColor(0xE8, 0xB8, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF4, 0xEF)
FOOTER_ON_NAVY = RGBColor(0xA8, 0xB4, 0xC4)


def style(run, text, size, bold, color, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def paint_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def title(slide, text, subtitle=None, top=0.85):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(8.6), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    style(p.add_run(), text, 32, True, NAVY)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        style(p2.add_run(), subtitle, 16, False, MUTED)


def bullets(slide, lines, top=2.4, size=18):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(8.6), Inches(4.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        style(p.add_run(), line, size, False, INK)


def footer(slide, text="barathx.com  ·  @getbarathx", color=MUTED):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(8.6), Inches(0.35))
    style(box.text_frame.paragraphs[0].add_run(), text, 12, False, color)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — Title
    s = prs.slides.add_slide(blank)
    paint_bg(s, NAVY)
    box = s.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(8.6), Inches(3.2))
    tf = box.text_frame
    tf.word_wrap = True
    style(tf.paragraphs[0].add_run(), "BarathX", 44, True, WHITE)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(14)
    style(p2.add_run(), "India’s public square — creator collab brief", 20, False, RGBColor(0xD8, 0xDE, 0xE6))
    p3 = tf.add_paragraph()
    p3.space_before = Pt(18)
    style(p3.add_run(), "Built by Indians · for Indians · owned by Indians", 14, False, SOFT)
    footer(s, color=FOOTER_ON_NAVY)

    # 2 — Problem
    s = prs.slides.add_slide(blank)
    paint_bg(s, BG)
    title(s, "The problem")
    bullets(
        s,
        [
            "Hot takes die in WhatsApp.",
            "Reels want your thumb — not your opinion.",
            "India needs a place to pick a side and argue it live.",
        ],
    )
    footer(s)

    # 3 — What
    s = prs.slides.add_slide(blank)
    paint_bg(s, BG)
    title(s, "What BarathX is")
    bullets(
        s,
        [
            "India’s public square — Square · Arenas · Live",
            "Drop a take → pick a side → real replies",
            "Human takes only. No AI slop.",
            "Built by Indians, for Indians, owned by Indians",
        ],
    )
    footer(s)

    # 4 — Why creators
    s = prs.slides.add_slide(blank)
    paint_bg(s, BG)
    title(s, "Why creators")
    bullets(
        s,
        [
            "Gen Z trusts niche creators more than ads.",
            "Your audience already argues in comments.",
            "We give them a square — you keep creative control.",
        ],
    )
    footer(s)

    # 5 — Ask
    s = prs.slides.add_slide(blank)
    paint_bg(s, BG)
    title(s, "The ask")
    bullets(
        s,
        [
            "1 Reel (or carousel) + 1 Story",
            "Exclusive BarathX link (we provide)",
            "CTA: create account → leave one take",
            "Spelling: BarathX",
        ],
    )
    footer(s)

    # 6 — Next
    s = prs.slides.add_slide(blank)
    paint_bg(s, BG)
    title(s, "Next step")
    bullets(
        s,
        [
            "Reply with how partnerships work on your side.",
            "We’ll share an exclusive link + short brief if it’s a fit.",
            "barathx.com",
            "@getbarathx",
        ],
    )
    footer(s)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
