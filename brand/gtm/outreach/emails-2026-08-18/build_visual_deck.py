#!/usr/bin/env python3
"""Presentable BarathX visual PPTX + one-pager PDF (live midnight screens)."""

from pathlib import Path
import shutil

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas as pdfcanvas

OUT_DIR = Path(__file__).resolve().parent
SLIDES_DIR = OUT_DIR / "_slide_assets"
LIVE = OUT_DIR / "screens-live"
FONTS = OUT_DIR / "_fonts"
PPTX_OUT = OUT_DIR / "BarathX-Creator-Collab-Brief.pptx"
PDF_OUT = OUT_DIR / "BarathX-One-Pager.pdf"
NIKHIL_PPTX = OUT_DIR / "BarathX-Nikhil-Brief.pptx"
NIKHIL_PDF = OUT_DIR / "BarathX-Nikhil-One-Pager.pdf"

W, H = 1920, 1080

BG = (13, 13, 18)
SURFACE = (22, 22, 30)
SURFACE2 = (28, 28, 38)
SAFFRON = (255, 153, 51)
SOFT = (255, 210, 160)
GREEN = (19, 136, 8)
WHITE = (255, 255, 255)
MUTED = (168, 172, 186)
LINE = (255, 153, 51, 55)

SCREENS = {
    "m_landing": LIVE / "m-landing.png",
    "m_login": LIVE / "m-login.png",
    "m_signup": LIVE / "m-signup.png",
    "d_landing": LIVE / "d-landing.png",
    "d_login": LIVE / "d-login.png",
    "ui_square": LIVE / "ui-square.jpg",
    "ui_arenas": LIVE / "ui-arenas.jpg",
}


def font_display(size):
    for p in [FONTS / "Syne.ttf", "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"]:
        if Path(p).exists():
            return ImageFont.truetype(str(p), size)
    return ImageFont.load_default()


def font_body(size, medium=False):
    prefer = [
        FONTS / "DMSans-Medium.ttf" if medium else FONTS / "DMSans-Regular.ttf",
        FONTS / "DMSans.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for p in prefer:
        if Path(p).exists():
            return ImageFont.truetype(str(p), size)
    return ImageFont.load_default()


def gradient_bg():
    img = Image.new("RGBA", (W, H), BG + (255,))
    glow = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    gd.ellipse((1150, -380, 2200, 580), fill=SAFFRON + (42,))
    gd.ellipse((-480, 560, 560, 1480), fill=(0, 0, 128, 36))
    gd.ellipse((720, 780, 1320, 1380), fill=GREEN + (16,))
    glow = glow.filter(ImageFilter.GaussianBlur(100))
    img = Image.alpha_composite(img, glow)
    # top saffron hairline
    d = ImageDraw.Draw(img)
    d.rectangle((0, 0, W, 4), fill=SAFFRON)
    return img


def round_resize(im: Image.Image, size, radius=40):
    im = im.convert("RGBA").resize(size, Image.Resampling.LANCZOS)
    mask = Image.new("L", size, 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, size[0] - 1, size[1] - 1), radius, fill=255)
    out = Image.new("RGBA", size, (0, 0, 0, 0))
    out.paste(im, (0, 0), mask)
    return out


def phone(screen_path: Path, height=820, radius=52, bezel=12):
    src = Image.open(screen_path).convert("RGBA")
    w0, h0 = src.size
    if h0 / w0 > 1.8:
        src = src.crop((0, 0, w0, int(w0 * 2.05)))
        w0, h0 = src.size
    ph = height
    pw = max(280, int(ph * (w0 / h0)))
    screen = round_resize(src, (pw - 2 * bezel, ph - 2 * bezel), radius=radius - 6)
    frame = Image.new("RGBA", (pw, ph), (0, 0, 0, 0))
    draw = ImageDraw.Draw(frame)
    shadow = Image.new("RGBA", (pw + 40, ph + 40), (0, 0, 0, 0))
    ImageDraw.Draw(shadow).rounded_rectangle((20, 24, pw + 20, ph + 24), radius + 4, fill=(0, 0, 0, 120))
    shadow = shadow.filter(ImageFilter.GaussianBlur(20))
    draw.rounded_rectangle((0, 0, pw - 1, ph - 1), radius, fill=(18, 18, 24, 255))
    draw.rounded_rectangle((1, 1, pw - 2, ph - 2), radius - 1, outline=SAFFRON + (90,), width=2)
    frame.paste(screen, (bezel, bezel), screen)
    nw, nh = int(pw * 0.30), 14
    draw.rounded_rectangle(((pw - nw) // 2, 10, (pw + nw) // 2, 10 + nh), 8, fill=(8, 8, 12, 255))
    return frame, shadow


def paste_phone(base, path, xy, height=820):
    frame, shadow = phone(path, height=height)
    sx, sy = xy
    base.alpha_composite(shadow, (sx - 20, sy - 16))
    base.alpha_composite(frame, (sx, sy))
    return frame.size


def paste_panel(base, path, box, radius=28):
    x, y, w, h = box
    im = round_resize(Image.open(path), (w, h), radius=radius)
    border = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    ImageDraw.Draw(border).rounded_rectangle((0, 0, w - 1, h - 1), radius, outline=SAFFRON + (90,), width=2)
    # soft shadow
    sh = Image.new("RGBA", (w + 30, h + 30), (0, 0, 0, 0))
    ImageDraw.Draw(sh).rounded_rectangle((12, 14, w + 12, h + 14), radius, fill=(0, 0, 0, 90))
    sh = sh.filter(ImageFilter.GaussianBlur(14))
    base.alpha_composite(sh, (x - 12, y - 10))
    base.alpha_composite(im, (x, y))
    base.alpha_composite(border, (x, y))


def t(draw, s, xy, size=48, display=True, fill=WHITE, medium=False, anchor=None):
    f = font_display(size) if display else font_body(size, medium=medium)
    draw.text(xy, s, font=f, fill=fill, anchor=anchor)


def pill(base, text, xy, fill=SAFFRON, fg=(13, 13, 18), pad_x=22, pad_y=12, size=16):
    f = font_body(size, medium=True)
    bbox = f.getbbox(text)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    w, h = tw + pad_x * 2, th + pad_y * 2
    layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    ImageDraw.Draw(layer).rounded_rectangle((0, 0, w - 1, h - 1), h // 2, fill=fill)
    ImageDraw.Draw(layer).text((pad_x, pad_y - 2), text, font=f, fill=fg)
    base.alpha_composite(layer, xy)
    return w, h


def card(base, xy, wh, title, body, accent=SAFFRON):
    x, y = xy
    w, h = wh
    layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(layer)
    d.rounded_rectangle((0, 0, w - 1, h - 1), 24, fill=SURFACE2 + (235,))
    d.rounded_rectangle((0, 0, w - 1, h - 1), 24, outline=accent + (70,), width=1)
    d.rectangle((0, 0, 8, h), fill=accent)
    td = ImageDraw.Draw(layer)
    td.text((28, 28), title, font=font_display(26), fill=WHITE)
    # wrap body
    f = font_body(18)
    words = body.split()
    lines, cur = [], ""
    max_w = w - 56
    for word in words:
        trial = f"{cur} {word}".strip()
        if f.getlength(trial) <= max_w:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    yy = 78
    for line in lines[:4]:
        td.text((28, yy), line, font=f, fill=MUTED)
        yy += 28
    base.alpha_composite(layer, (x, y))


def eyebrow(draw, label, xy=(110, 70)):
    t(draw, label.upper(), xy, 15, display=False, fill=SAFFRON, medium=True)


def footer(draw, text="barathx.com  ·  @getbarathx"):
    t(draw, text, (110, 1010), 18, display=False, fill=MUTED)


def slide_title():
    img = gradient_bg()
    d = ImageDraw.Draw(img)
    pill(img, "LIVE PRODUCT", (110, 90), fill=SAFFRON, fg=(13, 13, 18))
    t(d, "BarathX", (110, 200), 96)
    t(d, "India’s public square", (110, 320), 36, display=False, fill=SAFFRON, medium=True)
    t(d, "Built by Indians · for Indians · owned by Indians", (110, 385), 22, display=False, fill=SOFT)
    t(
        d,
        "Exactly what India needs from social —\nlive now, not a concept deck.",
        (110, 470),
        28,
        display=False,
        fill=MUTED,
    )
    footer(d)
    paste_phone(img, SCREENS["m_landing"], (1260, 80), height=920)
    return img


def slide_problem():
    img = gradient_bg()
    d = ImageDraw.Draw(img)
    eyebrow(d, "The problem")
    t(d, "India’s best takes\ndie in WhatsApp.", (110, 140), 60)
    card(img, (110, 420), (520, 220), "Foreign feeds", "Youth hours go to apps owned elsewhere. Algorithms reward rage, not discourse.")
    card(img, (660, 420), (520, 220), "No public square", "Comments get buried. Reels want your thumb — not your opinion.")
    footer(d)
    paste_phone(img, SCREENS["m_login"], (1280, 120), height=860)
    return img


def slide_product():
    img = gradient_bg()
    d = ImageDraw.Draw(img)
    eyebrow(d, "The product")
    t(d, "Drop a take. Pick a side.\nGet real replies.", (110, 120), 48)
    t(d, "Human takes only. No AI slop.", (110, 280), 22, display=False, fill=SAFFRON, medium=True)
    phones = [SCREENS["m_landing"], SCREENS["m_login"], SCREENS["m_signup"]]
    labels = ["Landing", "Sign in", "Create account"]
    hs = 700
    f, _ = phone(phones[0], height=hs)
    gap = 48
    total = f.width * 3 + gap * 2
    x0 = (W - total) // 2
    y0 = 340
    for i, p in enumerate(phones):
        size = paste_phone(img, p, (x0 + i * (f.width + gap), y0), height=hs)
        t(d, labels[i], (x0 + i * (f.width + gap) + size[0] // 2, y0 + size[1] + 24), 18, display=False, fill=MUTED, anchor="mm")
    return img


def slide_square():
    img = gradient_bg()
    d = ImageDraw.Draw(img)
    eyebrow(d, "Square · Arenas · Live")
    t(d, "India’s public square — live.", (110, 120), 48)
    paste_panel(img, SCREENS["ui_square"], (90, 230, 860, 740), radius=28)
    paste_panel(img, SCREENS["ui_arenas"], (990, 230, 860, 740), radius=28)
    return img


def slide_human():
    img = gradient_bg()
    d = ImageDraw.Draw(img)
    eyebrow(d, "Human first")
    t(d, "No AI slop.\nOnly real humans.", (110, 130), 56)
    card(img, (110, 380), (500, 200), "Flagged & demoted", "Likely-AI drafts are tagged. Real replies rise.")
    card(img, (640, 380), (500, 200), "Sided talk", "Agree / Disagree — stakes, not scroll theater.")
    card(img, (110, 620), (500, 200), "Live voices", "Live rooms: human voices only.")
    card(img, (640, 620), (500, 200), "On the record", "Takes don’t vanish in a group chat.")
    paste_phone(img, SCREENS["m_landing"], (1320, 160), height=800)
    return img


def slide_safety():
    img = gradient_bg()
    d = ImageDraw.Draw(img)
    eyebrow(d, "Safety we ship")
    t(d, "Safe square.\nNot an afterthought.", (110, 120), 52)
    card(img, (110, 360), (560, 210), "No adult content", "Blocked in posts, replies, DMs, Live. We do not encourage it.")
    card(img, (700, 360), (560, 210), "Report & guidelines", "In-product report paths. Community guidelines live.")
    card(img, (110, 610), (560, 210), "India DPDP", "Consent, export, delete. privacy@barathx.com")
    card(img, (700, 610), (560, 210), "Country blocks next", "Geo controls against abuse / adult traffic targeting us.")
    footer(d, "Trust is product — soft launch live")
    return img


def slide_no_monetization():
    img = gradient_bg()
    d = ImageDraw.Draw(img)
    eyebrow(d, "Incentives")
    t(d, "Not an ads feed.\nNot selling your mood.", (110, 130), 52)
    card(img, (110, 400), (520, 220), "Discourse first", "Built for sided debate — not engagement farming.")
    card(img, (660, 400), (520, 220), "No ads marketplace", "No pay-to-post. No cash-for-signup games.")
    footer(d)
    paste_panel(img, SCREENS["ui_square"], (1240, 180, 600, 720), radius=24)
    return img


def slide_india():
    img = gradient_bg()
    d = ImageDraw.Draw(img)
    eyebrow(d, "Built for India")
    t(d, "Exactly what you’re\nlooking for —\nspecially for India.", (110, 140), 48)
    t(d, "Built by Indians. For Indians. Owned by Indians.", (110, 420), 24, display=False, fill=SOFT)
    t(
        d,
        "Not a US clone with India skin.\nA public square for sided debate — made here.",
        (110, 500),
        22,
        display=False,
        fill=MUTED,
    )
    pill(img, "SIGN UP LIVE", (110, 620))
    footer(d)
    paste_panel(img, SCREENS["d_landing"], (980, 140, 860, 780), radius=24)
    return img


def slide_ask():
    img = gradient_bg()
    d = ImageDraw.Draw(img)
    eyebrow(d, "Live product · the ask")
    t(d, "Not an idea.\nA live site.", (110, 140), 56)
    t(
        d,
        "Sign up at barathx.com and explore —\nor we run a quick demo if you prefer.",
        (110, 340),
        24,
        display=False,
        fill=MUTED,
    )
    card(
        img,
        (110, 480),
        (900, 200),
        "GTM support — not build help",
        "Help more people discover and sign up at scale so we get real human validation from a large Indian audience.",
    )
    pill(img, "→  barathx.com  ·  sign up live", (110, 740), pad_x=28, pad_y=16, size=20)
    t(d, "@getbarathx  ·  hello@barathx.com", (110, 830), 18, display=False, fill=MUTED)
    paste_phone(img, SCREENS["m_landing"], (1280, 100), height=880)
    return img


def save_slide(img: Image.Image, name: str) -> Path:
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    path = SLIDES_DIR / f"{name}.png"
    img.convert("RGB").save(path, "PNG", optimize=True)
    return path


def build_pptx(paths):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for p in paths:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(PPTX_OUT)
    shutil.copy2(PPTX_OUT, NIKHIL_PPTX)
    print(f"Wrote {PPTX_OUT}")
    print(f"Wrote {NIKHIL_PPTX}")


def build_pdf():
    pw, ph = 1240, 1754
    canvas_img = Image.new("RGBA", (pw, ph), BG + (255,))
    glow = Image.new("RGBA", (pw, ph), (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    gd.ellipse((420, -280, 1380, 420), fill=SAFFRON + (48,))
    glow = glow.filter(ImageFilter.GaussianBlur(90))
    canvas_img = Image.alpha_composite(canvas_img, glow)
    d = ImageDraw.Draw(canvas_img)
    d.rectangle((0, 0, pw, 5), fill=SAFFRON)

    def tt(s, xy, size, display=True, fill=WHITE, medium=False):
        f = font_display(size) if display else font_body(size, medium=medium)
        d.text(xy, s, font=f, fill=fill)

    # LIVE pill
    pill(canvas_img, "LIVE", (64, 48), pad_x=16, pad_y=8, size=13)
    tt("BarathX", (64, 100), 54)
    tt("India’s public square", (64, 168), 24, display=False, fill=SAFFRON, medium=True)
    tt("Exactly what you’re looking for — specially for India.", (64, 210), 16, display=False, fill=SOFT)
    tt("Built by Indians · for Indians · owned by Indians", (64, 240), 14, display=False, fill=MUTED)

    phones = [SCREENS["m_landing"], SCREENS["m_login"], SCREENS["m_signup"]]
    f, _ = phone(phones[0], height=760, bezel=11, radius=46)
    gap = 28
    total = f.width * 3 + gap * 2
    x0 = (pw - total) // 2
    y0 = 290
    for i, p in enumerate(phones):
        fr, sh = phone(p, height=760, bezel=11, radius=46)
        canvas_img.alpha_composite(sh, (x0 + i * (f.width + gap) - 14, y0 - 10))
        canvas_img.alpha_composite(fr, (x0 + i * (f.width + gap), y0))

    y = y0 + f.height + 40
    tt("Drop a take. Pick a side. Get real replies.", (64, y), 24)
    chips = [
        "No AI slop",
        "No adult content",
        "Not an ads feed",
        "Indian-owned",
    ]
    cx, cy = 64, y + 48
    for chip in chips:
        w, h = pill(canvas_img, chip, (cx, cy), fill=SURFACE2, fg=SOFT, pad_x=16, pad_y=10, size=14)
        # outline
        ImageDraw.Draw(canvas_img).rounded_rectangle((cx, cy, cx + w - 1, cy + h - 1), h // 2, outline=SAFFRON + (100,))
        cx += w + 12

    # ask band
    band_y = ph - 280
    band = Image.new("RGBA", (pw - 128, 160), (0, 0, 0, 0))
    bd = ImageDraw.Draw(band)
    bd.rounded_rectangle((0, 0, pw - 129, 159), 28, fill=SURFACE2 + (240,))
    bd.rounded_rectangle((0, 0, pw - 129, 159), 28, outline=SAFFRON + (80,), width=1)
    canvas_img.alpha_composite(band, (64, band_y))
    tt("Not an idea — a live site.", (88, band_y + 28), 22)
    tt("Sign up & explore, or ask for a quick demo.", (88, band_y + 66), 16, display=False, fill=MUTED)
    tt("GTM support for large-scale signups → real human validation.", (88, band_y + 100), 15, display=False, fill=SOFT)

    cta_y = ph - 90
    pill(canvas_img, "→  barathx.com", (64, cta_y), pad_x=26, pad_y=14, size=18)
    tt("@getbarathx", (320, cta_y + 14), 16, display=False, fill=MUTED)

    tmp = SLIDES_DIR / "one-pager.png"
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    canvas_img.convert("RGB").save(tmp, "PNG", optimize=True)
    c = pdfcanvas.Canvas(str(PDF_OUT), pagesize=A4)
    page_w, page_h = A4
    c.drawImage(ImageReader(tmp), 0, 0, width=page_w, height=page_h, preserveAspectRatio=False, mask="auto")
    c.showPage()
    c.save()
    shutil.copy2(PDF_OUT, NIKHIL_PDF)
    print(f"Wrote {PDF_OUT}")
    print(f"Wrote {NIKHIL_PDF}")


def main():
    missing = [k for k, p in SCREENS.items() if not p.exists()]
    if missing:
        raise SystemExit(f"Missing live screens: {missing}")
    builders = [
        ("01-title", slide_title),
        ("02-problem", slide_problem),
        ("03-product", slide_product),
        ("04-square", slide_square),
        ("05-human", slide_human),
        ("06-safety", slide_safety),
        ("07-no-ads", slide_no_monetization),
        ("08-india", slide_india),
        ("09-ask", slide_ask),
    ]
    paths = []
    for name, fn in builders:
        print(f"render {name}…")
        paths.append(save_slide(fn(), name))
    build_pptx(paths)
    build_pdf()


if __name__ == "__main__":
    main()
